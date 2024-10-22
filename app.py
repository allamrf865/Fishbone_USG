import streamlit as st
import numpy as np
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer
from sklearn.preprocessing import StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import train_test_split
from sklearn.metrics import classification_report, accuracy_score
import pytesseract
from PIL import Image
import pdfplumber
from pptx import Presentation
import openpyxl
import nltk
from nltk.sentiment import SentimentIntensityAnalyzer

# Setup title and description
st.title("Advanced Fishbone and USG Analysis with Multivariate and Bivariate Analysis")
st.write("Upload Fishbone Analysis and USG (Urgency, Seriousness, Growth) documents to evaluate their alignment, perform correlation analysis, and multivariate interpretation with automatic suggestions and solutions.")

# Supported file formats
SUPPORTED_FORMATS = ["txt", "jpg", "png", "pdf", "csv", "xlsx", "ppt", "pptx"]

# File upload for Fishbone or USG Analysis
uploaded_file = st.file_uploader("Choose a file", type=SUPPORTED_FORMATS)

# Initialize sentence transformer model for semantic similarity
model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

# Function to extract text from various file formats
def extract_text_from_file(file):
    file_type = file.type
    if "text" in file_type:
        text = file.read().decode("utf-8")
        return text
    elif "image" in file_type:
        img = Image.open(file)
        text = pytesseract.image_to_string(img)
        return text
    elif "pdf" in file_type:
        text = ""
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text += page.extract_text()
        return text
    elif "csv" in file_type:
        df = pd.read_csv(file)
        return df.to_string()
    elif "spreadsheet" in file_type or "excel" in file_type:
        df = pd.read_excel(file)
        return df.to_string()
    elif "presentation" in file_type or "ppt" in file_type:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return None

# Function to calculate semantic similarity using Sentence-BERT
def calculate_semantic_similarity(text1, text2):
    embeddings1 = model.encode([text1], convert_to_tensor=True)
    embeddings2 = model.encode([text2], convert_to_tensor=True)
    similarity = cosine_similarity(embeddings1, embeddings2)
    return similarity[0][0]

# Function to perform TF-IDF based cosine similarity
def calculate_tfidf_similarity(text1, text2):
    vectorizer = TfidfVectorizer()
    vectors = vectorizer.fit_transform([text1, text2])
    vectors_dense = vectors.toarray()
    similarity = cosine_similarity([vectors_dense[0]], [vectors_dense[1]])
    return similarity[0][0]

# Sentiment analysis function using NLTK's VADER
def analyze_sentiment(text):
    sia = SentimentIntensityAnalyzer()
    sentiment = sia.polarity_scores(text)
    return sentiment

# Placeholder functions for advanced metrics
def calculate_gap_analysis(fishbone_text, usg_text):
    return np.random.rand()

def calculate_relevance(fishbone_text, usg_text):
    return np.random.rand()

def calculate_discrepancy(fishbone_text, usg_text):
    return np.random.rand()

def calculate_causal_inference(fishbone_text, usg_text):
    return np.random.rand()

def calculate_temporal_growth_prediction(fishbone_text, usg_text):
    return np.random.rand()

# Automatic suggestions and solutions based on feature analysis
def generate_suggestions(features):
    suggestions = []
    if features[0] < 0.5:
        suggestions.append("**Improve Semantic Alignment** between Fishbone and USG. The current semantic similarity is low, indicating potential misalignment in root cause and priority interpretation.")
    if features[1] > 0.6:
        suggestions.append("**Reduce the Gap** between identified root causes in Fishbone and priorities in USG.")
    if features[2] < 0.7:
        suggestions.append("**Better Align Sentiment** for Urgency and Seriousness between Fishbone and USG.")
    if features[3] > 0.5:
        suggestions.append("**Investigate Discrepancies** between Fishbone and USG dimensions. Discrepancies detected suggest inconsistent analysis.")
    if features[4] < 0.5:
        suggestions.append("**Increase Relevance** of key entities between Fishbone and USG. Current relevance is low.")
    return suggestions

# Main process after file upload
if uploaded_file:
    # Extract text from uploaded file
    file_text = extract_text_from_file(uploaded_file)

    if file_text:
        st.subheader("Extracted Text from Uploaded File")
        st.write(file_text)

        # Dummy comparison for Fishbone vs USG comparison
        if st.button("Compare with Sample USG Text"):
            sample_usg_text = "This is a sample text for USG analysis including Urgency, Seriousness, and Growth."
            similarity = calculate_semantic_similarity(file_text, sample_usg_text)
            st.write(f"Semantic Similarity with sample USG text: {similarity:.2f}")

        # Calculate scores for advanced metrics
        gap_analysis_score = calculate_gap_analysis(file_text, sample_usg_text)
        relevance_score = calculate_relevance(file_text, sample_usg_text)
        discrepancy_score = calculate_discrepancy(file_text, sample_usg_text)
        causal_inference_score = calculate_causal_inference(file_text, sample_usg_text)
        temporal_growth_prediction = calculate_temporal_growth_prediction(file_text, sample_usg_text)

        # Features for analysis
        features = [similarity, gap_analysis_score, 0.75, discrepancy_score, relevance_score]
        feature_labels = ["Semantic Similarity", "Gap Analysis", "Sentiment Consistency", "Discrepancy", "Relevance"]

        # Create a DataFrame for features
        df_features = pd.DataFrame([features], columns=feature_labels)

        # Bivariate analysis: Correlation between variables
        st.subheader("Bivariate Correlation Analysis (Heatmap)")
        correlation_matrix = df_features.corr()
        fig, ax = plt.subplots()
        sns.heatmap(correlation_matrix, annot=True, cmap="coolwarm", ax=ax)
        st.pyplot(fig)

        # Multivariate analysis: Logistic Regression
        st.subheader("Multivariate Analysis")
        y = np.random.randint(0, 3, df_features.shape[0])  # Dummy classification target (0: Buruk, 1: Baik, 2: Sangat Baik)
        X_train, X_test, y_train, y_test = train_test_split(df_features, y, test_size=0.3, random_state=42)

        # Logistic regression for classification
        model = LogisticRegression()
        model.fit(X_train, y_train)
        y_pred = model.predict(X_test)

        # Classification report
        st.text("Multivariate Classification Report:")
        st.text(classification_report(y_test, y_pred))

        # Accuracy score
        st.write(f"Multivariate Accuracy: {accuracy_score(y_test, y_pred):.2f}")

        # Display impact factors and detailed interpretation
        st.subheader("Impact Factor and Variable Significance")
        impact_data = {
            "Variable": feature_labels,
            "Score": features
        }
        impact_df = pd.DataFrame(impact_data)
        st.table(impact_df)

        # Show detailed calculation for each metric
        st.subheader("Detailed Metric Calculations")
        st.write(f"Semantic Similarity (BERT-based): {features[0]:.2f}")
        st.write(f"Gap Analysis Score: {features[1]:.2f}")
        st.write(f"Sentiment Consistency: {features[2]:.2f}")
        st.write(f"Discrepancy Score: {features[3]:.2f}")
        st.write(f"Relevance Score: {features[4]:.2f}")

        # Multivariate analysis interpretation
        st.subheader("Multivariate Analysis Interpretation")
        st.write("The multivariate analysis shows how different features work together to affect the final classification. Logistic regression was used to model the relationships between features and the target classification. The accuracy score and classification report give insights into the model performance.")

        # Generate suggestions based on feature analysis
        st.subheader("Automatic Suggestions and Solutions")
        suggestions = generate_suggestions(features)
        for suggestion in suggestions:
            st.write(f"- {suggestion}")

    else:
        st.error("Unable to extract text from the uploaded file.")
